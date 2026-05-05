"""
File Processor - Extract text from various file types

Handles PDF, Word, PowerPoint, text files, and more.
"""

import io
import os
from typing import Dict, Any, Optional
import logging

logger = logging.getLogger(__name__)


async def extract_text_from_file(file_content: bytes, filename: str, content_type: str) -> Dict[str, Any]:
    """
    Extract text from uploaded file

    Args:
        file_content: File content as bytes
        filename: Original filename
        content_type: MIME type

    Returns:
        Dictionary with:
        - text: Extracted text
        - metadata: File metadata (pages, word count, etc.)
        - asset_type: Type of file (pdf, docx, txt, etc.)
    """
    extension = os.path.splitext(filename)[1].lower()

    result = {
        "text": "",
        "metadata": {
            "filename": filename,
            "size": len(file_content),
            "content_type": content_type,
        },
        "asset_type": "file"
    }

    try:
        # PDF files
        if extension == ".pdf" or "pdf" in content_type.lower():
            result["text"] = await _extract_pdf(file_content)
            result["asset_type"] = "pdf"

        # Word documents
        elif extension in [".doc", ".docx"] or "word" in content_type.lower():
            result["text"] = await _extract_docx(file_content)
            result["asset_type"] = "docx"

        # PowerPoint
        elif extension in [".ppt", ".pptx"] or "powerpoint" in content_type.lower():
            result["text"] = await _extract_pptx(file_content)
            result["asset_type"] = "pptx"

        # Excel files
        elif extension in [".xlsx", ".xls"] or "excel" in content_type.lower() or "spreadsheet" in content_type.lower():
            result["text"] = await _extract_excel(file_content)
            result["asset_type"] = "xlsx"

        # Text files
        elif extension in [".txt", ".md", ".csv", ".json", ".xml", ".log"]:
            result["text"] = file_content.decode("utf-8", errors="ignore")
            result["asset_type"] = "text"

        # Code files
        elif extension in [".py", ".js", ".java", ".cpp", ".c", ".h", ".css", ".html", ".sql"]:
            result["text"] = file_content.decode("utf-8", errors="ignore")
            result["asset_type"] = "code"

        else:
            # Try to decode as text
            try:
                result["text"] = file_content.decode("utf-8", errors="strict")
                result["asset_type"] = "text"
            except UnicodeDecodeError:
                result["text"] = f"File type '{extension}' not supported for text extraction"
                result["metadata"]["error"] = "Unsupported file type"

    except Exception as e:
        logger.error(f"Failed to extract text from {filename}: {e}")
        result["text"] = f"Error extracting text: {str(e)}"
        result["metadata"]["error"] = str(e)

    # Add metadata
    result["metadata"]["word_count"] = len(result["text"].split())
    result["metadata"]["char_count"] = len(result["text"])

    return result


async def _extract_pdf(content: bytes) -> str:
    """Extract text from PDF file"""
    try:
        # Try PyPDF2 first (most reliable)
        try:
            import PyPDF2
            from io import BytesIO

            pdf_file = BytesIO(content)
            reader = PyPDF2.PdfReader(pdf_file)

            text_parts = []
            for page_num in range(len(reader.pages)):
                page = reader.pages[page_num]
                text_parts.append(page.extract_text())

            return "\n\n".join(text_parts)

        except ImportError:
            logger.warning("PyPDF2 not installed, falling back to pdfplumber")

        # Fallback to pdfplumber
        try:
            import pdfplumber
            from io import BytesIO

            pdf_file = BytesIO(content)
            text_parts = []

            with pdfplumber.open(pdf_file) as pdf:
                for page in pdf.pages:
                    text = page.extract_text()
                    if text:
                        text_parts.append(text)

            return "\n\n".join(text_parts)

        except ImportError:
            return "PDF extraction requires PyPDF2 or pdfplumber. Install with: pip install PyPDF2"

    except Exception as e:
        logger.error(f"PDF extraction failed: {e}")
        return f"Failed to extract PDF text: {str(e)}"


async def _extract_docx(content: bytes) -> str:
    """Extract text from Word document"""
    try:
        import docx
        from io import BytesIO

        doc_file = BytesIO(content)
        doc = docx.Document(doc_file)

        text_parts = []

        # Extract paragraphs
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text)

        # Extract tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text for cell in row.cells)
                if row_text.strip():
                    text_parts.append(row_text)

        return "\n\n".join(text_parts)

    except ImportError:
        return "Word document extraction requires python-docx. Install with: pip install python-docx"
    except Exception as e:
        logger.error(f"DOCX extraction failed: {e}")
        return f"Failed to extract Word document text: {str(e)}"


async def _extract_pptx(content: bytes) -> str:
    """Extract text from PowerPoint presentation"""
    try:
        from pptx import Presentation
        from io import BytesIO

        ppt_file = BytesIO(content)
        prs = Presentation(ppt_file)

        text_parts = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = f"=== Slide {slide_num} ===\n"

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text += shape.text + "\n"

            if slide_text.strip() != f"=== Slide {slide_num} ===":
                text_parts.append(slide_text)

        return "\n\n".join(text_parts)

    except ImportError:
        return "PowerPoint extraction requires python-pptx. Install with: pip install python-pptx"
    except Exception as e:
        logger.error(f"PPTX extraction failed: {e}")
        return f"Failed to extract PowerPoint text: {str(e)}"


async def _extract_excel(content: bytes) -> str:
    """Extract text from Excel spreadsheet"""
    try:
        import openpyxl
        from io import BytesIO

        excel_file = BytesIO(content)
        workbook = openpyxl.load_workbook(excel_file, data_only=True)

        text_parts = []

        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]

            # Add sheet header
            text_parts.append(f"=== Sheet: {sheet_name} ===\n")

            # Extract all rows
            rows_text = []
            for row in sheet.iter_rows(values_only=True):
                # Filter out completely empty rows
                if any(cell is not None and str(cell).strip() for cell in row):
                    # Convert cells to strings, handling None values
                    row_text = " | ".join(
                        str(cell) if cell is not None else ""
                        for cell in row
                    )
                    rows_text.append(row_text)

            if rows_text:
                text_parts.append("\n".join(rows_text))

        return "\n\n".join(text_parts)

    except ImportError:
        return "Excel extraction requires openpyxl. Install with: pip install openpyxl"
    except Exception as e:
        logger.error(f"Excel extraction failed: {e}")
        return f"Failed to extract Excel text: {str(e)}"

