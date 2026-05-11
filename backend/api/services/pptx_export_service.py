"""
PPTX Export Service

Converts structured slide data to PowerPoint (.pptx) format using python-pptx library.
Handles different slide layouts and applies themes (colors, fonts).
"""

import io
import os
import tempfile
from typing import List, Dict, Any, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import logging
import httpx
from PIL import Image
import matplotlib
matplotlib.use('Agg')  # Use non-interactive backend
import matplotlib.pyplot as plt

logger = logging.getLogger(__name__)


class PPTXExportService:
    """Service for exporting presentations to PowerPoint format"""

    def __init__(self):
        self.slide_width = Inches(10)
        self.slide_height = Inches(7.5)
        self.temp_dir = tempfile.gettempdir()

    async def _download_image(self, url: str) -> Optional[str]:
        """
        Download image from URL and save to temp file.

        Args:
            url: Image URL

        Returns:
            Path to downloaded image file, or None if download fails
        """
        try:
            async with httpx.AsyncClient(timeout=30.0) as client:
                response = await client.get(url)
                response.raise_for_status()

                # Validate content type
                content_type = response.headers.get("content-type", "")
                if not content_type.startswith("image/"):
                    logger.warning(f"URL {url} is not an image (content-type: {content_type})")
                    return None

                # Check file size (max 5MB)
                content_length = len(response.content)
                if content_length > 5 * 1024 * 1024:
                    logger.warning(f"Image {url} is too large ({content_length} bytes)")
                    return None

                # Validate it's a valid image using Pillow
                try:
                    img = Image.open(io.BytesIO(response.content))
                    img.verify()  # Verify it's a valid image
                except Exception as e:
                    logger.warning(f"Invalid image from {url}: {e}")
                    return None

                # Save to temp file
                extension = content_type.split("/")[-1]
                if extension not in ["png", "jpeg", "jpg", "gif", "bmp"]:
                    extension = "png"

                temp_path = os.path.join(
                    self.temp_dir,
                    f"slide_image_{hash(url)}_{os.getpid()}.{extension}"
                )

                with open(temp_path, "wb") as f:
                    f.write(response.content)

                logger.info(f"Downloaded image from {url} to {temp_path}")
                return temp_path

        except httpx.HTTPError as e:
            logger.error(f"Failed to download image from {url}: {e}")
            return None
        except Exception as e:
            logger.error(f"Unexpected error downloading image from {url}: {e}")
            return None

    def _generate_chart(
        self,
        chart_data: Dict[str, Any],
        chart_type: str,
        theme: Dict[str, Any]
    ) -> Optional[str]:
        """
        Generate chart image using matplotlib.

        Args:
            chart_data: Chart data with labels, values, title, etc.
            chart_type: Type of chart (bar, line, pie, scatter)
            theme: Theme configuration for colors

        Returns:
            Path to generated chart image, or None if generation fails
        """
        try:
            # Extract data
            labels = chart_data.get("labels", [])
            values = chart_data.get("values", [])
            title = chart_data.get("title", "")
            x_label = chart_data.get("x_label", "")
            y_label = chart_data.get("y_label", "")

            if not values:
                logger.warning("No chart values provided")
                return None

            # Get theme colors
            primary_color = theme.get("colors", {}).get("primary", "#0066cc")

            # Create figure
            fig, ax = plt.subplots(figsize=(8, 6))

            # Generate chart based on type
            if chart_type == "bar":
                ax.bar(labels if labels else range(len(values)), values, color=primary_color)
            elif chart_type == "line":
                ax.plot(labels if labels else range(len(values)), values, color=primary_color, marker='o')
            elif chart_type == "pie":
                ax.pie(values, labels=labels, autopct='%1.1f%%', colors=[primary_color])
            elif chart_type == "scatter":
                # For scatter, values should be a list of [x, y] pairs
                if isinstance(values[0], (list, tuple)) and len(values[0]) == 2:
                    x_vals = [v[0] for v in values]
                    y_vals = [v[1] for v in values]
                    ax.scatter(x_vals, y_vals, color=primary_color, s=100)
                else:
                    # Fallback: use index as x
                    ax.scatter(range(len(values)), values, color=primary_color, s=100)
            else:
                logger.warning(f"Unsupported chart type: {chart_type}")
                return None

            # Set labels
            if title:
                ax.set_title(title, fontsize=16, fontweight='bold')
            if x_label and chart_type != "pie":
                ax.set_xlabel(x_label, fontsize=12)
            if y_label and chart_type != "pie":
                ax.set_ylabel(y_label, fontsize=12)

            # Grid for non-pie charts
            if chart_type != "pie":
                ax.grid(True, alpha=0.3)

            # Save to temp file
            temp_path = os.path.join(
                self.temp_dir,
                f"slide_chart_{os.getpid()}_{id(chart_data)}.png"
            )

            plt.tight_layout()
            fig.savefig(temp_path, dpi=150, bbox_inches='tight')
            plt.close(fig)

            logger.info(f"Generated {chart_type} chart at {temp_path}")
            return temp_path

        except Exception as e:
            logger.error(f"Failed to generate chart: {e}")
            return None

    def export_to_pptx(
        self,
        slides: List[Dict[str, Any]],
        theme: Dict[str, Any],
        title: str = "Presentation"
    ) -> bytes:
        """
        Convert structured slide data to PPTX file.

        Args:
            slides: List of slide dictionaries with content_json and slide_type
            theme: Theme configuration with colors and fonts
            title: Presentation title

        Returns:
            PPTX file as bytes
        """
        try:
            prs = Presentation()
            prs.slide_width = self.slide_width
            prs.slide_height = self.slide_height

            for slide_data in slides:
                slide_type = slide_data.get("slide_type")
                content = slide_data.get("content_json", {})
                speaker_notes = slide_data.get("speaker_notes", "")

                if slide_type == "title":
                    self._add_title_slide(prs, content, theme, speaker_notes)
                elif slide_type == "bullets":
                    self._add_bullet_slide(prs, content, theme, speaker_notes)
                elif slide_type == "two_column":
                    self._add_two_column_slide(prs, content, theme, speaker_notes)
                elif slide_type == "content":
                    self._add_content_slide(prs, content, theme, speaker_notes)
                elif slide_type == "image_text":
                    self._add_image_text_slide(prs, content, theme, speaker_notes)
                elif slide_type == "chart":
                    self._add_chart_slide(prs, content, theme, speaker_notes)
                else:
                    logger.warning(f"Unknown slide type: {slide_type}, using content layout")
                    self._add_content_slide(prs, content, theme, speaker_notes)

            # Save to bytes
            pptx_io = io.BytesIO()
            prs.save(pptx_io)
            pptx_io.seek(0)
            return pptx_io.getvalue()

        except Exception as e:
            logger.error(f"Failed to export PPTX: {str(e)}")
            raise

    def _add_title_slide(
        self,
        prs: Presentation,
        content: Dict[str, Any],
        theme: Dict[str, Any],
        speaker_notes: str = ""
    ):
        """Create elegant title slide with background design"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Use blank layout for full control

        # Add sophisticated background design
        self._add_background_design(slide, theme, style="gradient")

        # Add title with enhanced styling
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5),
            Inches(8), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = content.get("title", "")
        title_frame.word_wrap = True

        # Center align title
        for paragraph in title_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER

        # Apply theme with larger title font
        colors = theme.get("colors", {})
        fonts = theme.get("fonts", {})
        font_name = fonts.get("heading", "Calibri")
        color_hex = colors.get("primary", "#000000")

        for paragraph in title_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = font_name
                run.font.size = Pt(54)  # Extra large for title slide
                run.font.bold = True
                if color_hex:
                    rgb = self._hex_to_rgb(color_hex)
                    run.font.color.rgb = RGBColor(*rgb)

        # Add subtitle with elegant positioning
        subtitle_text = content.get("subtitle", "")
        if not subtitle_text and content.get("elements"):
            subtitle_text = content["elements"][0].get("content", "")

        if subtitle_text:
            subtitle_box = slide.shapes.add_textbox(
                Inches(1.5), Inches(4.2),
                Inches(7), Inches(1)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle_text
            subtitle_frame.word_wrap = True

            # Center align subtitle
            for paragraph in subtitle_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER

            # Apply theme to subtitle
            subtitle_color = colors.get("secondary", colors.get("text", "#666666"))
            for paragraph in subtitle_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = font_name
                    run.font.size = Pt(24)
                    if subtitle_color:
                        rgb = self._hex_to_rgb(subtitle_color)
                        run.font.color.rgb = RGBColor(*rgb)

        # Add decorative accent line
        accent_color = colors.get("accent", colors.get("primary", "#0066cc"))
        accent_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(3.5), Inches(5.5),
            Inches(3), Inches(0.08)
        )
        from pptx.enum.shapes import MSO_SHAPE
        accent_line.fill.solid()
        rgb = self._hex_to_rgb(accent_color)
        accent_line.fill.fore_color.rgb = RGBColor(*rgb)
        accent_line.line.fill.background()

        # Add speaker notes
        if speaker_notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = speaker_notes

    def _add_bullet_slide(
        self,
        prs: Presentation,
        content: Dict[str, Any],
        theme: Dict[str, Any],
        speaker_notes: str = ""
    ):
        """Create elegant bullet point slide with background design"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout for control

        # Add background design
        self._add_background_design(slide, theme, style="corner_accent")

        # Add title with enhanced styling
        title_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(0.5),
            Inches(8.6), Inches(0.9)
        )
        title_frame = title_box.text_frame
        title_frame.text = content.get("title", "")
        self._apply_theme_to_text(title_box, theme, "heading")

        # Add decorative underline for title
        colors = theme.get("colors", {})
        accent_color = colors.get("accent", colors.get("primary", "#0066cc"))
        title_underline = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.7), Inches(1.35),
            Inches(1.5), Inches(0.05)
        )
        from pptx.enum.shapes import MSO_SHAPE
        title_underline.fill.solid()
        rgb = self._hex_to_rgb(accent_color)
        title_underline.fill.fore_color.rgb = RGBColor(*rgb)
        title_underline.line.fill.background()

        # Add bullet points with improved layout
        body_box = slide.shapes.add_textbox(
            Inches(0.9), Inches(1.8),
            Inches(8.2), Inches(5)
        )
        tf = body_box.text_frame
        tf.word_wrap = True
        tf.clear()

        elements = content.get("elements", [])
        for idx, item in enumerate(elements):
            if item.get("type") == "bullet":
                if idx == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()

                p.text = item.get("content", "")
                p.level = item.get("level", 0)

                # Enhanced spacing
                p.space_before = Pt(12)
                p.space_after = Pt(12)
                p.line_spacing = 1.3

        self._apply_theme_to_text(body_box, theme, "body")

        # Add speaker notes
        if speaker_notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = speaker_notes

    def _add_two_column_slide(
        self,
        prs: Presentation,
        content: Dict[str, Any],
        theme: Dict[str, Any],
        speaker_notes: str = ""
    ):
        """Create elegant two-column slide with visual separation"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        # Add background design
        self._add_background_design(slide, theme, style="bottom_accent")

        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(0.5),
            Inches(8.6), Inches(0.9)
        )
        title_frame = title_box.text_frame
        title_frame.text = content.get("title", "")
        self._apply_theme_to_text(title_box, theme, "heading")

        # Add decorative underline for title
        colors = theme.get("colors", {})
        accent_color = colors.get("accent", colors.get("primary", "#0066cc"))
        from pptx.enum.shapes import MSO_SHAPE
        title_underline = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.7), Inches(1.35),
            Inches(1.5), Inches(0.05)
        )
        title_underline.fill.solid()
        rgb = self._hex_to_rgb(accent_color)
        title_underline.fill.fore_color.rgb = RGBColor(*rgb)
        title_underline.line.fill.background()

        # Add vertical separator line between columns
        separator = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(5), Inches(1.8),
            Inches(0.03), Inches(5)
        )
        separator.fill.solid()
        sep_rgb = self._hex_to_rgb(accent_color)
        separator.fill.fore_color.rgb = RGBColor(sep_rgb[0], sep_rgb[1], sep_rgb[2])
        separator.fill.fore_color.brightness = 0.5  # Lighter separator
        separator.line.fill.background()

        elements = content.get("elements", [])

        # Separate left and right column content
        left_content = [e for e in elements if e.get("column") == "left"]
        right_content = [e for e in elements if e.get("column") == "right"]

        # Add left column with enhanced styling
        if left_content:
            left_box = slide.shapes.add_textbox(
                Inches(0.7), Inches(1.8),
                Inches(4), Inches(5)
            )
            left_frame = left_box.text_frame
            left_frame.word_wrap = True

            for idx, item in enumerate(left_content):
                if idx > 0:
                    p = left_frame.add_paragraph()
                else:
                    p = left_frame.paragraphs[0]
                p.text = "• " + item.get("content", "")
                p.level = 0
                p.space_before = Pt(8)
                p.space_after = Pt(8)
                p.line_spacing = 1.2

            self._apply_theme_to_text(left_box, theme, "body")

        # Add right column with enhanced styling
        if right_content:
            right_box = slide.shapes.add_textbox(
                Inches(5.3), Inches(1.8),
                Inches(4), Inches(5)
            )
            right_frame = right_box.text_frame
            right_frame.word_wrap = True

            for idx, item in enumerate(right_content):
                if idx > 0:
                    p = right_frame.add_paragraph()
                else:
                    p = right_frame.paragraphs[0]
                p.text = "• " + item.get("content", "")
                p.level = 0
                p.space_before = Pt(8)
                p.space_after = Pt(8)
                p.line_spacing = 1.2

            self._apply_theme_to_text(right_box, theme, "body")

        # Add speaker notes
        if speaker_notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = speaker_notes

    def _add_content_slide(
        self,
        prs: Presentation,
        content: Dict[str, Any],
        theme: Dict[str, Any],
        speaker_notes: str = ""
    ):
        """Create elegant content slide with enhanced styling"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        # Add background design
        self._add_background_design(slide, theme, style="accent_bar")

        # Add title with enhanced styling
        title_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(0.5),
            Inches(8.6), Inches(0.9)
        )
        title_frame = title_box.text_frame
        title_frame.text = content.get("title", "")
        self._apply_theme_to_text(title_box, theme, "heading")

        # Add decorative underline for title
        colors = theme.get("colors", {})
        accent_color = colors.get("accent", colors.get("primary", "#0066cc"))
        from pptx.enum.shapes import MSO_SHAPE
        title_underline = slide.shapes.add_textbox(
            MSO_SHAPE.RECTANGLE,
            Inches(0.7), Inches(1.35),
            Inches(1.5), Inches(0.05)
        )
        title_underline.fill.solid()
        rgb = self._hex_to_rgb(accent_color)
        title_underline.fill.fore_color.rgb = RGBColor(*rgb)
        title_underline.line.fill.background()

        # Add content with improved readability
        body_box = slide.shapes.add_textbox(
            Inches(0.9), Inches(1.8),
            Inches(8.2), Inches(5)
        )
        tf = body_box.text_frame
        tf.word_wrap = True
        tf.clear()

        elements = content.get("elements", [])
        for idx, item in enumerate(elements):
            if idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.text = "• " + item.get("content", "")
            p.level = 0

            # Enhanced spacing for readability
            p.space_before = Pt(10)
            p.space_after = Pt(10)
            p.line_spacing = 1.3

        self._apply_theme_to_text(body_box, theme, "body")

        # Add speaker notes
        if speaker_notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = speaker_notes

    async def _add_image_text_slide(
        self,
        prs: Presentation,
        content: Dict[str, Any],
        theme: Dict[str, Any],
        speaker_notes: str = ""
    ):
        """Create slide with image and text"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5),
            Inches(9), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = content.get("title", "")
        self._apply_theme_to_text(title_box, theme, "heading")

        # Get image and text elements
        elements = content.get("elements", [])
        image_elements = [e for e in elements if e.get("type") == "image"]
        text_elements = [e for e in elements if e.get("type") != "image"]

        # Add text content (left side)
        if text_elements:
            text_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.5),
                Inches(4.5), Inches(5.5)
            )
            text_frame = text_box.text_frame
            text_frame.word_wrap = True

            for item in text_elements:
                p = text_frame.add_paragraph()
                p.text = item.get("content", "")

            self._apply_theme_to_text(text_box, theme, "body")

        # Add image (right side) - actual image insertion
        if image_elements:
            image_url = image_elements[0].get("content", "")
            if image_url:
                # This will be called async, so we need to handle it properly
                # For now, add a note about async requirement
                import asyncio
                try:
                    # Try to download image
                    image_path = asyncio.run(self._download_image(image_url))
                    if image_path and os.path.exists(image_path):
                        # Add actual image
                        slide.shapes.add_picture(
                            image_path,
                            Inches(5.25), Inches(1.5),
                            width=Inches(4.25), height=Inches(5.5)
                        )
                        # Clean up temp file
                        try:
                            os.remove(image_path)
                        except:
                            pass
                    else:
                        # Fallback to placeholder if download failed
                        self._add_image_placeholder(slide, theme)
                except Exception as e:
                    logger.warning(f"Failed to add image from {image_url}: {e}")
                    # Fallback to placeholder
                    self._add_image_placeholder(slide, theme)
            else:
                self._add_image_placeholder(slide, theme)
        else:
            self._add_image_placeholder(slide, theme)

        # Add speaker notes
        if speaker_notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = speaker_notes

    def _add_image_placeholder(self, slide, theme):
        """Add a placeholder box for images"""
        image_box = slide.shapes.add_textbox(
            Inches(5.25), Inches(1.5),
            Inches(4.25), Inches(5.5)
        )
        image_frame = image_box.text_frame
        image_frame.text = "[Image Placeholder]"
        self._apply_theme_to_text(image_box, theme, "body")

    def _add_chart_slide(
        self,
        prs: Presentation,
        content: Dict[str, Any],
        theme: Dict[str, Any],
        speaker_notes: str = ""
    ):
        """Create slide with actual chart"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5),
            Inches(9), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = content.get("title", "")
        self._apply_theme_to_text(title_box, theme, "heading")

        # Extract chart data and type
        elements = content.get("elements", [])
        chart_data = content.get("chart_data", {})
        chart_type = content.get("chart_type", "bar")

        # If chart_data is in elements, extract it
        if not chart_data and elements:
            for element in elements:
                if element.get("type") == "chart":
                    chart_data = element.get("chart_data", {})
                    chart_type = element.get("chart_type", "bar")
                    break

        # Generate chart if we have data
        if chart_data and chart_data.get("values"):
            try:
                chart_path = self._generate_chart(chart_data, chart_type, theme)
                if chart_path and os.path.exists(chart_path):
                    # Add chart image to slide (centered)
                    slide.shapes.add_picture(
                        chart_path,
                        Inches(1), Inches(1.5),
                        width=Inches(8), height=Inches(5.5)
                    )
                    # Clean up temp file
                    try:
                        os.remove(chart_path)
                    except:
                        pass
                else:
                    # Fallback to placeholder
                    self._add_chart_placeholder(slide, theme)
            except Exception as e:
                logger.warning(f"Failed to generate chart: {e}")
                self._add_chart_placeholder(slide, theme)
        else:
            # No chart data, add placeholder
            self._add_chart_placeholder(slide, theme)

        # Add speaker notes
        if speaker_notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = speaker_notes

    def _add_chart_placeholder(self, slide, theme):
        """Add a placeholder for charts"""
        chart_box = slide.shapes.add_textbox(
            Inches(1), Inches(2),
            Inches(8), Inches(4)
        )
        chart_frame = chart_box.text_frame
        chart_frame.text = "[Chart Placeholder]\n\nChart data not available"
        self._apply_theme_to_text(chart_box, theme, "body")

        # Add speaker notes
        if speaker_notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = speaker_notes

    def _apply_theme_to_text(
        self,
        shape,
        theme: Dict[str, Any],
        text_type: str
    ):
        """
        Apply theme colors and fonts to text shape with enhanced styling.

        Args:
            shape: PowerPoint shape with text
            theme: Theme configuration
            text_type: "heading" or "body"
        """
        try:
            colors = theme.get("colors", {})
            fonts = theme.get("fonts", {})

            font_name = fonts.get(text_type, "Arial")
            color_key = "primary" if text_type == "heading" else "text"
            color_hex = colors.get(color_key, "#000000")

            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = font_name
                    if color_hex:
                        rgb = self._hex_to_rgb(color_hex)
                        run.font.color.rgb = RGBColor(*rgb)

                    # Enhanced font sizing and styling
                    if text_type == "heading":
                        run.font.size = Pt(40)  # Larger, more impactful
                        run.font.bold = True
                    else:
                        run.font.size = Pt(20)  # Slightly larger for readability

                # Add spacing between paragraphs
                if text_type == "body":
                    paragraph.space_before = Pt(6)
                    paragraph.space_after = Pt(6)

        except Exception as e:
            logger.warning(f"Failed to apply theme: {str(e)}")

    def _add_background_design(self, slide, theme: Dict[str, Any], style: str = "gradient"):
        """
        Add elegant background design to slide based on template theme.

        Args:
            slide: PowerPoint slide object
            theme: Theme configuration with colors
            style: Background style - "gradient", "accent_bar", "corner_accent", "minimal"
        """
        try:
            from pptx.enum.shapes import MSO_SHAPE
            colors = theme.get("colors", {})
            primary_color = colors.get("primary", "#0066cc")
            secondary_color = colors.get("secondary", "#00aaff")
            accent_color = colors.get("accent", "#ff6600")
            background_color = colors.get("background", "#ffffff")

            # Add subtle background if not white
            if background_color and background_color.lower() != "#ffffff":
                bg_shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    0, 0,
                    self.slide_width, self.slide_height
                )
                bg_shape.fill.solid()
                rgb = self._hex_to_rgb(background_color)
                bg_shape.fill.fore_color.rgb = RGBColor(*rgb)
                bg_shape.line.fill.background()
                # Send to back
                slide.shapes._spTree.remove(bg_shape._element)
                slide.shapes._spTree.insert(2, bg_shape._element)

            if style == "gradient" or style == "modern":
                # Add gradient accent shape in corner
                accent_shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    self.slide_width - Inches(3),
                    self.slide_height - Inches(2.5),
                    Inches(3.5), Inches(2.5)
                )
                accent_shape.fill.solid()
                rgb = self._hex_to_rgb(primary_color)
                accent_shape.fill.fore_color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])
                # Make semi-transparent
                accent_shape.fill.fore_color.brightness = 0.8
                accent_shape.line.fill.background()
                # Rotate for dynamic look
                accent_shape.rotation = -15
                # Send to back
                slide.shapes._spTree.remove(accent_shape._element)
                slide.shapes._spTree.insert(2, accent_shape._element)

            elif style == "accent_bar":
                # Add vertical accent bar on left
                bar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    0, 0,
                    Inches(0.3), self.slide_height
                )
                bar.fill.solid()
                rgb = self._hex_to_rgb(accent_color)
                bar.fill.fore_color.rgb = RGBColor(*rgb)
                bar.line.fill.background()
                # Send to back
                slide.shapes._spTree.remove(bar._element)
                slide.shapes._spTree.insert(2, bar._element)

            elif style == "corner_accent":
                # Add accent triangle in corner
                triangle = slide.shapes.add_shape(
                    MSO_SHAPE.RIGHT_TRIANGLE,
                    self.slide_width - Inches(2), 0,
                    Inches(2), Inches(2)
                )
                triangle.fill.solid()
                rgb = self._hex_to_rgb(secondary_color)
                triangle.fill.fore_color.rgb = RGBColor(*rgb)
                triangle.line.fill.background()
                triangle.rotation = 180
                # Send to back
                slide.shapes._spTree.remove(triangle._element)
                slide.shapes._spTree.insert(2, triangle._element)

            elif style == "bottom_accent":
                # Add decorative bottom bar
                bar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    0, self.slide_height - Inches(0.2),
                    self.slide_width, Inches(0.2)
                )
                bar.fill.solid()
                rgb = self._hex_to_rgb(primary_color)
                bar.fill.fore_color.rgb = RGBColor(*rgb)
                bar.line.fill.background()

        except Exception as e:
            logger.warning(f"Failed to add background design: {str(e)}")

    def _add_decorative_elements(self, slide, theme: Dict[str, Any], position: str = "header"):
        """
        Add decorative design elements to enhance slide aesthetics.

        Args:
            slide: PowerPoint slide object
            theme: Theme configuration
            position: Where to add elements - "header", "footer", "corner"
        """
        try:
            from pptx.enum.shapes import MSO_SHAPE
            colors = theme.get("colors", {})
            accent_color = colors.get("accent", "#ff6600")

            if position == "header":
                # Add small accent line under title area
                line = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(0.5), Inches(1.2),
                    Inches(2), Inches(0.05)
                )
                line.fill.solid()
                rgb = self._hex_to_rgb(accent_color)
                line.fill.fore_color.rgb = RGBColor(*rgb)
                line.line.fill.background()

            elif position == "corner":
                # Add small decorative circles
                for i, offset in enumerate([0, 0.3, 0.6]):
                    circle = slide.shapes.add_shape(
                        MSO_SHAPE.OVAL,
                        Inches(9 + offset), Inches(0.2),
                        Inches(0.15), Inches(0.15)
                    )
                    circle.fill.solid()
                    rgb = self._hex_to_rgb(accent_color)
                    # Vary opacity
                    circle.fill.fore_color.rgb = RGBColor(*rgb)
                    circle.line.fill.background()

        except Exception as e:
            logger.warning(f"Failed to add decorative elements: {str(e)}")

    def _hex_to_rgb(self, hex_color: str) -> tuple:
        """Convert hex color to RGB tuple"""
        hex_color = hex_color.lstrip('#')
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
